import asyncio

import pptx

from app.processors.base import BaseProcessor


class PPTProcessor(BaseProcessor):
    async def extract(self, file_path: str) -> list[str]:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._extract_sync, file_path)

    def _extract_sync(self, file_path: str) -> list[str]:
        prs = pptx.Presentation(file_path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return slides if slides else []
